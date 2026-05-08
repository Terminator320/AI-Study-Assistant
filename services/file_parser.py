import os
import io

MAX_CHARS = int(os.getenv("MAX_FILE_TEXT_CHARS", 12000))


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        return _from_pdf(file_bytes)
    elif ext == "pptx":
        return _from_pptx(file_bytes)
    else:
        return file_bytes.decode("utf-8", errors="ignore")[:MAX_CHARS]


def _from_pdf(data: bytes) -> str:
    import PyPDF2
    reader = PyPDF2.PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)[:MAX_CHARS]


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
    return "\n".join(parts)[:MAX_CHARS]
